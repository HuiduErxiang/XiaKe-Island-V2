"""Test fixtures and generators for input-structure Skill."""

import io
import os
import sys
import tempfile
from pathlib import Path

import pytest

# Ensure the skill package is importable
sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent))


@pytest.fixture(autouse=True)
def _auto_generate_fixtures(fixtures_dir):
    """Automatically generate test fixture files before any test runs."""
    _ensure_fixtures(fixtures_dir)


@pytest.fixture
def fixtures_dir():
    return Path(__file__).resolve().parent / "fixtures"


@pytest.fixture
def output_dir():
    with tempfile.TemporaryDirectory() as d:
        yield Path(d)


def _ensure_fixtures(fixtures_dir: Path):
    """Generate test fixture files if they don't exist."""
    fixtures_dir.mkdir(parents=True, exist_ok=True)

    # --- Text PDF (3 pages) ---
    text_pdf = fixtures_dir / "sample_text.pdf"
    if not text_pdf.exists():
        _create_text_pdf(text_pdf)

    # --- Table PDF (2 pages) ---
    table_pdf = fixtures_dir / "sample_table.pdf"
    if not table_pdf.exists():
        _create_table_pdf(table_pdf)

    # --- PDF with chart image ---
    chart_pdf = fixtures_dir / "sample_chart.pdf"
    if not chart_pdf.exists():
        _create_chart_pdf(chart_pdf)

    # --- PPTX (5 slides) ---
    pptx_file = fixtures_dir / "sample_presentation.pptx"
    if not pptx_file.exists():
        _create_pptx(pptx_file)

    # --- DOCX (multi-level headings) ---
    docx_file = fixtures_dir / "sample_report.docx"
    if not docx_file.exists():
        _create_docx(docx_file)

    # --- TXT (UTF-8 paragraphs) ---
    txt_file = fixtures_dir / "sample_notes.txt"
    if not txt_file.exists():
        _create_txt(txt_file)

    # --- TXT (GBK encoding) ---
    gbk_file = fixtures_dir / "sample_gbk.txt"
    if not gbk_file.exists():
        _create_gbk_txt(gbk_file)

    # --- Empty file ---
    empty_file = fixtures_dir / "empty.txt"
    if not empty_file.exists():
        empty_file.write_text("")


def _create_text_pdf(path: Path):
    """Create a 3-page text PDF with title hierarchy."""
    import fitz  # PyMuPDF

    doc = fitz.open()
    # Page 1
    page = doc.new_page()
    page.insert_text((72, 72), "Introduction to Lecanemab", fontsize=18, fontname="helv")
    page.insert_text((72, 110), "This document provides an overview of lecanemab,", fontsize=11, fontname="helv")
    page.insert_text((72, 126), "a novel therapeutic agent for Alzheimer's disease.", fontsize=11, fontname="helv")
    page.insert_text((72, 150), "Author: van Dyck CH", fontsize=10, fontname="helv")
    page.insert_text((72, 166), "Source: New England Journal of Medicine, 2023", fontsize=10, fontname="helv")

    # Page 2
    page = doc.new_page()
    page.insert_text((72, 72), "Clinical Trial Design", fontsize=16, fontname="helv")
    page.insert_text((72, 110), "Participants were randomly assigned to receive", fontsize=11, fontname="helv")
    page.insert_text((72, 126), "intravenous lecanemab or placebo every 2 weeks.", fontsize=11, fontname="helv")
    page.insert_text((72, 150), "The primary endpoint was the change from baseline", fontsize=11, fontname="helv")
    page.insert_text((72, 166), "in the CDR-SB score at 18 months.", fontsize=11, fontname="helv")

    # Page 3
    page = doc.new_page()
    page.insert_text((72, 72), "Results and Conclusions", fontsize=16, fontname="helv")
    page.insert_text((72, 110), "The study met its primary endpoint, demonstrating", fontsize=11, fontname="helv")
    page.insert_text((72, 126), "a statistically significant reduction in clinical", fontsize=11, fontname="helv")
    page.insert_text((72, 142), "decline compared to placebo (p<0.001).", fontsize=11, fontname="helv")

    doc.save(str(path))
    doc.close()


def _create_table_pdf(path: Path):
    """Create a 2-page PDF with a simple table."""
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Efficacy Results", fontsize=16, fontname="helv")
    # Simulate a table with tab-separated text
    header = "Indicator\tLecanemab (n=898)\tPlacebo (n=897)\tp-value"
    row1 = "CDR-SB change\t1.21\t1.66\t<0.001"
    row2 = "Amyloid burden\t-55.48\t3.64\t<0.001"
    row3 = "ADAS-Cog14\t-1.44\t-0.13\t<0.001"
    page.insert_text((72, 110), header, fontsize=10, fontname="helv")
    page.insert_text((72, 128), row1, fontsize=10, fontname="helv")
    page.insert_text((72, 146), row2, fontsize=10, fontname="helv")
    page.insert_text((72, 164), row3, fontsize=10, fontname="helv")

    page = doc.new_page()
    page.insert_text((72, 72), "Safety Data", fontsize=16, fontname="helv")
    page.insert_text((72, 110), "Adverse Event\tLecanemab\tPlacebo", fontsize=10, fontname="helv")
    page.insert_text((72, 128), "ARIA-E\t12.6%\t1.7%", fontsize=10, fontname="helv")
    page.insert_text((72, 146), "Headache\t11.1%\t8.1%", fontsize=10, fontname="helv")

    doc.save(str(path))
    doc.close()


def _create_chart_pdf(path: Path):
    """Create a 2-page PDF with an embedded image (simulating a chart)."""
    import fitz
    from PIL import Image, ImageDraw

    # Create a simple chart image
    img = Image.new("RGB", (400, 300), "white")
    draw = ImageDraw.Draw(img)
    draw.rectangle((50, 100, 100, 250), fill="blue", outline="black")
    draw.rectangle((150, 80, 200, 250), fill="red", outline="black")
    draw.rectangle((250, 150, 300, 250), fill="green", outline="black")
    draw.text((60, 255), "Q1", fill="white")
    draw.text((160, 255), "Q2", fill="white")
    draw.text((260, 255), "Q3", fill="white")

    img_bytes = io.BytesIO()
    img.save(img_bytes, format="PNG")
    img_bytes.seek(0)

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "CDR-SB Score Over Time", fontsize=16, fontname="helv")
    page.insert_text((72, 100), "Figure 1: CDR-SB change from baseline", fontsize=10, fontname="helv")
    # Embed chart image
    rect = fitz.Rect(72, 120, 272, 320)
    page.insert_image(rect, stream=img_bytes.read())

    page = doc.new_page()
    page.insert_text((72, 72), "Additional Analysis", fontsize=16, fontname="helv")
    page.insert_text((72, 110), "Supporting data for the primary endpoint.", fontsize=11, fontname="helv")

    doc.save(str(path))
    doc.close()


def _create_pptx(path: Path):
    """Create a 5-slide PPTX."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for i in range(1, 6):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = f"Slide {i} — Product Overview Part {i}"
        body = slide.placeholders[1]
        body.text = f"Key point for slide {i}:\n• Feature A description\n• Feature B comparison"
        notes = slide.notes_slide if slide.has_notes_slide else None
        if notes is None:
            notes = slide.notes_slide
        notes.notes_text_frame.text = f"Speaker notes for slide {i}"

    prs.save(str(path))


def _create_docx(path: Path):
    """Create a DOCX with multi-level headings."""
    from docx import Document

    doc = Document()
    doc.add_heading("Clinical Research Report", level=1)
    doc.add_paragraph("This report summarizes the key findings from the phase III clinical trial.")

    doc.add_heading("Background", level=2)
    doc.add_paragraph("Alzheimer's disease affects approximately 50 million people worldwide.")

    doc.add_heading("Methods", level=2)
    doc.add_paragraph("A randomized, double-blind, placebo-controlled trial was conducted.")

    doc.add_heading("Participant Selection", level=3)
    doc.add_paragraph("Participants aged 50-90 with early Alzheimer's disease were eligible.")

    doc.add_heading("Statistical Analysis", level=3)
    doc.add_paragraph("The primary analysis used a mixed-effects model for repeated measures.")

    doc.add_heading("Results", level=2)
    doc.add_paragraph("The study enrolled 1795 participants across 235 sites in North America, Europe, and Asia.")

    doc.save(str(path))


def _create_txt(path: Path):
    """Create a UTF-8 TXT with paragraphs."""
    content = (
        "Clinical Trial Notes\n"
        "\n"
        "The lecanemab trial demonstrated significant cognitive benefits.\n"
        "Patients receiving the drug showed 27% less decline on CDR-SB.\n"
        "\n"
        "Safety concerns were primarily related to ARIA-E events.\n"
        "These were mostly asymptomatic and resolved within 4-6 weeks.\n"
        "\n"
        "The FDA granted accelerated approval in January 2023.\n"
        "Full approval followed in July 2023 after confirmatory data.\n"
    )
    path.write_text(content, encoding="utf-8")


def _create_gbk_txt(path: Path):
    """Create a GBK-encoded TXT."""
    content = "这是一份中文临床笔记。\n\n患者对治疗反应良好。\n不良反应观察中。\n"
    path.write_bytes(content.encode("gbk"))
