"""PPTX parser: extract slides, notes, and images."""

import logging
import os
from pathlib import Path

from ..schema import MaterialFormat, ParsedMaterial, ParseError
from ..utils.image_extract import extract_images_from_pptx
from .. import config

logger = logging.getLogger(__name__)


def parse_pptx(path: Path, output_dir: Path) -> ParsedMaterial:
    """Parse a PPTX file into structured Markdown.

    Args:
        path: Path to the PPTX file.
        output_dir: Output directory for extracted assets.

    Returns:
        ParsedMaterial with slide-by-slide Markdown content.
    """
    from pptx import Presentation

    filename = os.path.basename(str(path))

    file_size = os.path.getsize(str(path))
    if file_size > config.MAX_FILE_SIZE_BYTES:
        return ParsedMaterial(
            filename=filename,
            format=MaterialFormat.PPTX,
            content="",
            error=ParseError(filename=filename, reason="file_too_large"),
        )

    try:
        prs = Presentation(str(path))
    except Exception as e:
        return ParsedMaterial(
            filename=filename,
            format=MaterialFormat.PPTX,
            content="",
            error=ParseError(filename=filename, reason="parse_error"),
        )

    slides_md: list[str] = []
    all_images = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        try:
            title = _get_slide_title(slide)
            body = _get_slide_body(slide)
            notes = _get_slide_notes(slide)

            title_text = title or f"Slide {slide_num}"
            md = f"### Slide {slide_num} — {title_text}\n\n"

            if body:
                md += f"{body}\n\n"

            if notes:
                md += f"> **备注**: {notes}\n"

            slides_md.append(md)

            # Extract images
            images = extract_images_from_pptx(
                slide, output_dir / "assets", filename, slide_num
            )
            for img in images:
                if not img.caption:
                    img.caption = f"图片 — Slide {slide_num}"
            all_images.extend(images)

        except Exception as e:
            logger.warning(f"Error processing slide {slide_num}: {e}")
            slides_md.append(f"### Slide {slide_num}\n\n[Slide {slide_num} — 解析错误]\n")

    return ParsedMaterial(
        filename=filename,
        format=MaterialFormat.PPTX,
        content="\n\n".join(slides_md),
        page_count=len(prs.slides),
        images=all_images,
    )


def _get_slide_title(slide) -> str:
    """Extract the title text from a slide's title placeholder."""
    if slide.shapes.title and slide.shapes.title.text:
        return slide.shapes.title.text.strip()
    # Fallback: first text shape that looks like a title
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            first_line = shape.text_frame.text.strip().split("\n")[0]
            if len(first_line) < 100:
                return first_line
    return ""


def _get_slide_body(slide) -> str:
    """Extract body text from non-title shapes."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    body_parts = []
    title_shape = slide.shapes.title
    for shape in slide.shapes:
        if shape == title_shape:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                body_parts.append(text)
    return "\n".join(body_parts)


def _get_slide_notes(slide) -> str:
    """Extract speaker notes."""
    try:
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            return notes
    except Exception:
        pass
    return ""
